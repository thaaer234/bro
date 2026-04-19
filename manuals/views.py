from io import BytesIO
import base64
from datetime import datetime
from pathlib import Path
import subprocess
import tempfile

import qrcode
from PIL import UnidentifiedImageError
from django.core import signing
from django.urls import NoReverseMatch
from django.urls import reverse
from django.contrib.auth.mixins import LoginRequiredMixin
from django.contrib.auth.models import User
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from django.templatetags.static import static
from django.views.generic import TemplateView
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Mm, Pt

from .guide_data import USER_MANUAL_ERRORS, build_manuals_context, build_user_manual_context


MANUAL_EXPORT_SIGNING_SALT = "manuals.user.handbook.export"


def _qr_data_uri(value):
    qr = qrcode.QRCode(box_size=7, border=2)
    qr.add_data(value)
    qr.make(fit=True)
    image = qr.make_image(fill_color="black", back_color="white")
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    encoded = base64.b64encode(buffer.getvalue()).decode("ascii")
    return f"data:image/png;base64,{encoded}"


def _build_closing_page_context(request):
    base_url = request.build_absolute_uri("/").rstrip("/")
    manual_url = request.build_absolute_uri("/manuals/handbook/")
    developer_site_url = "https://thaaer7426.space.z.ai/"
    return {
        "closing_brand_name": "نظام معهد اليمان",
        "closing_developer_name": "ثائر المصري",
        "closing_role": "Full-Stack Developer",
        "closing_stack": "Django / Web Systems",
        "closing_developer_phone": "0983232446",
        "closing_developer_email": "thaaer74@gmail.com",
        "closing_developer_site_url": developer_site_url,
        "closing_developer_site_qr": _qr_data_uri(developer_site_url),
        "closing_developer_photo": request.build_absolute_uri("/media/profile_pictures/IMG_4472_2.jpeg"),
        "closing_skills": ["البرمجة", "تطوير الويب", "قواعد البيانات"],
        "closing_site_url": base_url,
        "closing_manual_url": manual_url,
        "closing_site_qr": _qr_data_uri(base_url),
        "closing_manual_qr": _qr_data_uri(manual_url),
    }


def _manual_target_users(request_user):
    if request_user.is_superuser:
        users = User.objects.select_related("employee_profile").order_by("first_name", "last_name", "username")
    else:
        users = User.objects.filter(pk=request_user.pk).select_related("employee_profile")

    items = []
    for user in users:
        employee = getattr(user, "employee_profile", None)
        items.append(
            {
                "id": user.pk,
                "username": user.username,
                "full_name": user.get_full_name() or user.username,
                "position": employee.get_position_display() if employee else "مستخدم نظام",
            }
        )
    return items


def _selected_manual_user(request):
    selected_id = request.GET.get("user") or request.POST.get("user")
    users = _manual_target_users(request.user)
    if not users:
        return None, users
    if selected_id:
        for item in users:
            if str(item["id"]) == str(selected_id):
                return get_object_or_404(User, pk=item["id"]), users
    return get_object_or_404(User, pk=users[0]["id"]), users


def _build_user_manual_identity_context(request, target_user):
    employee = getattr(target_user, "employee_profile", None)
    system_url = "https://alyaman-institute.com"
    password_value = (request.GET.get("password") or request.POST.get("password") or "").strip()
    return {
        "manual_target_user": target_user,
        "manual_target_username": target_user.username,
        "manual_target_name": target_user.get_full_name() or target_user.username,
        "manual_target_position": employee.get_position_display() if employee else "مستخدم نظام",
        "manual_target_phone": getattr(employee, "phone_number", "") or "غير محدد",
        "manual_target_email": target_user.email or "غير محدد",
        "manual_target_password": password_value or "غير متاح في النظام",
        "manual_target_qr": _qr_data_uri(system_url),
        "manual_target_system_url": system_url,
        "manual_target_system_code": "ALYAMAN-INSTITUTE",
    }


def _build_manual_error_pages():
    error_visuals = {
        "خطأ عدم الصلاحية": static("img/manual-center/error-permission-screen.svg"),
        "انتهاء الجلسة أو تسجيل الخروج التلقائي": static("img/manual-center/error-session-screen.svg"),
        "فشل التحقق من الحقول": static("img/manual-center/error-validation-screen.svg"),
        "العنصر غير موجود أو تم حذفه": static("img/manual-center/error-missing-screen.svg"),
        "تعارض أو تكرار بيانات": static("img/manual-center/error-conflict-screen.svg"),
        "فشل الشبكة أو بطء التحميل": static("img/manual-center/error-network-screen.svg"),
        "فشل الطباعة أو التصدير": static("img/manual-center/error-print-screen.svg"),
        "خطأ داخلي غير متوقع": static("img/manual-center/error-internal-screen.svg"),
    }
    pages = []
    for item in USER_MANUAL_ERRORS:
        screenshot_path = error_visuals.get(item["title"], static("img/manual-center/error-internal-screen.svg"))
        pages.append(
            {
                "title": item["title"],
                "group": "الأخطاء والتصرف الصحيح",
                "goal": "هذه الصفحة تشرح هذا الخطأ بطريقة تشغيلية: متى يظهر، كيف تكتشف سببه، وما هو التصرف الصحيح قبل تصعيده للدعم.",
                "used_by": "كل المستخدمين بحسب الصلاحيات",
                "path": "قد يظهر في أكثر من صفحة حسب نوع العملية",
                "screenshot": {
                    "title": item["title"],
                    "path": screenshot_path,
                    "caption": "لقطة مرجعية لهذا النوع من رسائل الخطأ داخل النظام حتى تتعرف على شكله العام قبل التعامل معه.",
                },
                "buttons": [
                    {
                        "label": "متى يظهر هذا الخطأ",
                        "location": "أثناء تنفيذ العملية الحالية أو عند فتح صفحة لا تستوفي الشروط المطلوبة.",
                        "used_by": "أي مستخدم قد يمر بنفس الحالة.",
                        "when": "فور ظهور الرسالة أو فشل الحفظ أو إعادة التوجيه غير المتوقعة.",
                        "purpose": item["when"],
                        "result": "تفهم سبب ظهور الخطأ قبل إعادة المحاولة بشكل عشوائي.",
                    },
                    {
                        "label": "التحقق الأولي",
                        "location": "في نفس الصفحة وقبل تكرار العملية.",
                        "used_by": "المستخدم الذي ظهرت له المشكلة.",
                        "when": "مباشرة بعد ظهور الخطأ.",
                        "purpose": "راجع الحقول، الصلاحيات، السجل المستهدف، والاتصال قبل أي إعادة تنفيذ.",
                        "result": "تحدد بسرعة هل المشكلة من البيانات أو الصلاحيات أو الشبكة أو من السجل نفسه.",
                    },
                    {
                        "label": "طريقة التصرف الصحيحة",
                        "location": "الخطوة التالية بعد فهم السبب.",
                        "used_by": "صاحب العملية الحالية.",
                        "when": "بعد التحقق الأولي مباشرة.",
                        "purpose": item["action"],
                        "result": "تعالج الخطأ بالطريقة الصحيحة أو تمنع تكراره عند المتابعة.",
                    },
                    {
                        "label": "متى ترفع للدعم",
                        "location": "بعد فشل المعالجة المباشرة أو تكرر الخطأ.",
                        "used_by": "أي مستخدم لا يستطيع إكمال العمل بعد المحاولة الصحيحة.",
                        "when": "إذا تكرر الخطأ بعد إعادة التنفيذ مرة واحدة فقط.",
                        "purpose": "جهز اسم الصفحة واسم العملية ووقت الخطأ ثم ارفعها للدعم أو للإدارة التقنية.",
                        "result": "يصل البلاغ واضحًا وقابلًا للمعالجة السريعة دون ضياع التفاصيل.",
                    },
                ],
            }
        )
    return pages


def _get_manual_param(params, key, default=""):
    value = params.get(key, default) if params is not None else default
    if isinstance(value, (list, tuple)):
        return value[0] if value else default
    return value


def _resolve_manual_target_user_from_params(request, params):
    selected_id = (_get_manual_param(params, "user", "") or "").strip()
    if selected_id:
        return get_object_or_404(User, pk=selected_id)
    selected_user, _users = _selected_manual_user(request)
    return selected_user


def _chunk_items(items, size):
    if size <= 0:
        return [items]
    return [items[index:index + size] for index in range(0, len(items), size)]


def _continuation_title(title, chunk_index, chunk_total):
    return title


def _text_weight(value, base=1, ratio=220):
    text = str(value or "").strip()
    return base + max(0, len(text) // ratio)


def _chunk_weighted(items, max_weight, weight_fn):
    if max_weight <= 0:
        return [items]
    chunks = []
    current = []
    current_weight = 0
    for item in items:
        item_weight = max(1, weight_fn(item))
        if current and (current_weight + item_weight) > max_weight:
            chunks.append(current)
            current = [item]
            current_weight = item_weight
        else:
            current.append(item)
            current_weight += item_weight
    if current:
        chunks.append(current)
    return chunks


def _chunk_weighted_progressive(items, first_weight, continuation_weight, weight_fn):
    if not items:
        return []
    first_chunk = _chunk_weighted(items, first_weight, weight_fn)
    if len(first_chunk) <= 1:
        return first_chunk

    chunks = [first_chunk[0]]
    remaining = items[len(first_chunk[0]):]
    if remaining:
        chunks.extend(_chunk_weighted(remaining, continuation_weight, weight_fn))
    return chunks


def _workflow_step_weight(step):
    return _text_weight(step, base=3, ratio=140)


def _action_weight(item):
    return (
        _text_weight(item.get("label"), base=1, ratio=80)
        + _text_weight(item.get("location"), base=2, ratio=120)
        + _text_weight(item.get("used_by"), base=1, ratio=120)
        + _text_weight(item.get("when"), base=2, ratio=120)
        + _text_weight(item.get("purpose"), base=2, ratio=120)
        + _text_weight(item.get("result"), base=2, ratio=120)
    )


def _page_status_label(chunk_index, chunk_total):
    if chunk_total <= 1:
        return "مكتملة"
    if chunk_index == 1:
        return "بداية القسم"
    if chunk_index == chunk_total:
        return "ختام القسم"
    return "متابعة"


def _should_use_visual_page(page_type, items):
    item_count = len(items or [])
    total_weight = sum(_action_weight(item) for item in items) if items and page_type != "workflow" else 0
    if page_type == "workflow":
        return item_count >= 5
    return item_count >= 4 or total_weight >= 38


def _build_user_handbook_content_pages(manual_workflows, manual_screens, manual_error_pages):
    pages = []

    for workflow_index, workflow in enumerate(manual_workflows, start=1):
        step_chunks = _chunk_weighted_progressive(workflow["steps"], 20, 26, _workflow_step_weight)
        step_counter = 1
        for chunk_index, steps in enumerate(step_chunks, start=1):
            chunk_items = [
                {"index": step_counter + offset, "text": step}
                for offset, step in enumerate(steps)
            ]
            use_visual_page = chunk_index == 1 and _should_use_visual_page("workflow", chunk_items)
            if use_visual_page:
                pages.append(
                    {
                        "type": "workflow",
                        "layout": "visual",
                        "title": workflow["title"],
                        "toc_title": workflow["title"],
                        "toc_kind": "مرحلة",
                        "toc_include": True,
                        "eyebrow": f"مرحلة {workflow_index}",
                        "intro": workflow["intro"],
                        "image_title": workflow["screenshot"]["title"],
                        "image_path": workflow["screenshot"]["path"],
                        "image_caption": workflow["screenshot"]["caption"],
                        "show_image": True,
                        "items": [],
                        "continued": False,
                    }
                )
            pages.append(
                {
                    "type": "workflow",
                    "layout": "content",
                    "title": _continuation_title(workflow["title"], chunk_index, len(step_chunks)),
                    "toc_title": workflow["title"],
                    "toc_kind": "مرحلة",
                    "toc_include": chunk_index == 1 and not use_visual_page,
                    "eyebrow": f"مرحلة {workflow_index}",
                    "intro": (
                        "استكمال الخطوات التنفيذية لنفس المرحلة مع الحفاظ على نفس التسلسل العملي."
                        if use_visual_page or chunk_index > 1
                        else workflow["intro"]
                    ),
                    "image_title": workflow["screenshot"]["title"],
                    "image_path": workflow["screenshot"]["path"],
                    "image_caption": workflow["screenshot"]["caption"],
                    "show_image": chunk_index == 1 and not use_visual_page,
                    "image_mode": "hero" if chunk_index == 1 and len(steps) <= 4 else "full",
                    "items": chunk_items,
                    "continued": chunk_index > 1 or use_visual_page,
                }
            )
            step_counter += len(steps)

    for screen in manual_screens:
        button_chunks = _chunk_weighted_progressive(screen["buttons"], 46, 58, _action_weight)
        button_counter = 1
        for chunk_index, buttons in enumerate(button_chunks, start=1):
            chunk_items = [
                {
                    **button,
                    "index": button_counter + offset,
                }
                for offset, button in enumerate(buttons)
            ]
            use_visual_page = chunk_index == 1 and _should_use_visual_page("screen", chunk_items)
            if use_visual_page:
                pages.append(
                    {
                        "type": "screen",
                        "layout": "visual",
                        "title": screen["title"],
                        "toc_title": screen["title"],
                        "toc_kind": screen["group"],
                        "toc_include": True,
                        "eyebrow": screen["group"],
                        "intro": screen["goal"],
                        "used_by": screen["used_by"],
                        "path": screen["path"],
                        "image_title": screen["screenshot"]["title"],
                        "image_path": screen["screenshot"]["path"],
                        "image_caption": screen["screenshot"]["caption"],
                        "show_image": True,
                        "items": [],
                        "continued": False,
                    }
                )
            pages.append(
                {
                    "type": "screen",
                    "layout": "content",
                    "title": _continuation_title(screen["title"], chunk_index, len(button_chunks)),
                    "toc_title": screen["title"],
                    "toc_kind": screen["group"],
                    "toc_include": chunk_index == 1 and not use_visual_page,
                    "eyebrow": screen["group"],
                    "intro": (
                        "متابعة عناصر الشاشة نفسها بعد تخصيص صفحة مستقلة لعرض الشاشة."
                        if use_visual_page
                        else screen["goal"] if chunk_index == 1
                        else "متابعة عناصر الشاشة نفسها بعد توزيع المحتوى على صفحة إضافية."
                    ),
                    "used_by": screen["used_by"],
                    "path": screen["path"],
                    "image_title": screen["screenshot"]["title"],
                    "image_path": screen["screenshot"]["path"],
                    "image_caption": screen["screenshot"]["caption"],
                    "show_image": chunk_index == 1 and not use_visual_page,
                    "image_mode": "hero" if chunk_index == 1 and len(buttons) <= 4 else "full",
                    "items": chunk_items,
                    "continued": chunk_index > 1 or use_visual_page,
                }
            )
            button_counter += len(buttons)

    for error_page in manual_error_pages:
        button_chunks = _chunk_weighted_progressive(error_page["buttons"], 42, 54, _action_weight)
        button_counter = 1
        for chunk_index, buttons in enumerate(button_chunks, start=1):
            chunk_items = [
                {
                    **button,
                    "index": button_counter + offset,
                }
                for offset, button in enumerate(buttons)
            ]
            use_visual_page = chunk_index == 1 and _should_use_visual_page("error", chunk_items)
            if use_visual_page:
                pages.append(
                    {
                        "type": "error",
                        "layout": "visual",
                        "title": error_page["title"],
                        "toc_title": error_page["title"],
                        "toc_kind": "خطأ",
                        "toc_include": True,
                        "eyebrow": error_page["group"],
                        "intro": error_page["goal"],
                        "used_by": error_page["used_by"],
                        "path": error_page["path"],
                        "image_title": error_page["screenshot"]["title"],
                        "image_path": error_page["screenshot"]["path"],
                        "image_caption": error_page["screenshot"]["caption"],
                        "show_image": True,
                        "items": [],
                        "continued": False,
                    }
                )
            pages.append(
                {
                    "type": "error",
                    "layout": "content",
                    "title": _continuation_title(error_page["title"], chunk_index, len(button_chunks)),
                    "toc_title": error_page["title"],
                    "toc_kind": "خطأ",
                    "toc_include": chunk_index == 1 and not use_visual_page,
                    "eyebrow": error_page["group"],
                    "intro": (
                        "استكمال عناصر التعامل مع الخطأ نفسه بعد تخصيص صفحة مستقلة للرسالة المرجعية."
                        if use_visual_page
                        else error_page["goal"] if chunk_index == 1
                        else "استكمال عناصر التعامل مع الخطأ نفسه في صفحة إضافية أوضح للطباعة."
                    ),
                    "used_by": error_page["used_by"],
                    "path": error_page["path"],
                    "image_title": error_page["screenshot"]["title"],
                    "image_path": error_page["screenshot"]["path"],
                    "image_caption": error_page["screenshot"]["caption"],
                    "show_image": chunk_index == 1 and not use_visual_page,
                    "image_mode": "hero" if chunk_index == 1 and len(buttons) <= 4 else "full",
                    "items": chunk_items,
                    "continued": chunk_index > 1 or use_visual_page,
                }
            )
            button_counter += len(buttons)

    return pages


def _build_user_handbook_toc_pages(content_pages):
    toc_page_size = 22
    toc_entries = [
        {"title": "مقدمة الدليل", "kind": "تمهيد", "page": 1},
        {"title": "بطاقة المستخدم والوصول", "kind": "وصول", "page": 2},
        {"title": "الفهرس", "kind": "تنقل", "page": 3},
    ]
    content_entries = [{"title": "بطاقة الوصول وبيانات المستخدم", "kind": "وصول"}]
    content_entries.extend(
        {
            "title": page["toc_title"],
            "kind": page["toc_kind"],
        }
        for page in content_pages
        if page.get("toc_include")
    )

    toc_pages = _chunk_items(toc_entries, toc_page_size)
    while True:
        recalculated_entries = list(toc_entries)
        page_number = 3 + len(toc_pages)

        for page in content_pages:
            if page.get("toc_include"):
                recalculated_entries.append(
                    {
                        "title": page["toc_title"],
                        "kind": page["toc_kind"],
                        "page": page_number,
                    }
                )
            page_number += 1

        recalculated_entries.extend(
            [
                {"title": "ملاحظات وتشغيل آمن", "kind": "تنبيهات", "page": page_number},
                {"title": "الأسئلة الشائعة", "kind": "FAQ", "page": page_number + 1},
                {"title": "خاتمة الدليل", "kind": "ختام", "page": page_number + 2},
            ]
        )
        new_toc_pages = _chunk_items(recalculated_entries, toc_page_size)
        if len(new_toc_pages) == len(toc_pages):
            return [
                {
                    "entries": page_entries,
                    "toc_page_number": 3 + index,
                }
                for index, page_entries in enumerate(new_toc_pages)
            ]
        toc_pages = new_toc_pages


def _build_user_handbook_page_context(request, params=None, *, print_mode=False):
    active_params = params if params is not None else request.GET
    target_user = _resolve_manual_target_user_from_params(request, active_params)
    query = (_get_manual_param(active_params, "q", "") or "").strip().lower()
    selected_group = (_get_manual_param(active_params, "group", "") or "").strip()

    context = {
        "manual_selected_user": target_user,
    }
    context.update(
        build_user_manual_context(
            target_user,
            query=query,
            selected_group=selected_group,
        )
    )
    context.update(_build_user_manual_identity_context(request, target_user))
    context["manual_error_pages"] = _build_manual_error_pages()
    context["manual_content_pages"] = _build_user_handbook_content_pages(
        context["manual_workflows"],
        context["manual_screens"],
        context["manual_error_pages"],
    )
    context["manual_toc_pages"] = _build_user_handbook_toc_pages(context["manual_content_pages"])
    context["manual_print_mode"] = print_mode
    context["manual_print_url"] = request.build_absolute_uri()
    context["manual_now_date"] = datetime.now().strftime("%Y / %m / %d")
    export_page = (_get_manual_param(active_params, "export_page", "") or "").strip()
    context["manual_export_page_index"] = int(export_page) if export_page.isdigit() else None
    context["manual_export_mode"] = (_get_manual_param(active_params, "export_mode", "") or "").strip()
    context["manual_total_render_pages"] = _manual_total_render_pages(context)
    context.update(_build_closing_page_context(request))
    return context


def _manual_total_render_pages(context):
    return 8 + len(context["manual_toc_pages"]) + len(context["manual_content_pages"])


def _find_edge_executable():
    candidates = [
        Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
        Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def _reverse_manual_url(*names):
    for name in names:
        try:
            return reverse(name)
        except NoReverseMatch:
            continue
    raise NoReverseMatch(", ".join(names))


def _pptx_set_text(paragraph, text, *, size=18, bold=False, color=(30, 41, 55), align=PP_ALIGN.RIGHT):
    paragraph.text = str(text or "")
    paragraph.alignment = align
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor(*color)


def _pptx_add_panel(slide, left, top, width, height, *, fill=(255, 255, 255), line=(217, 226, 236)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*line)
    return shape


def _pptx_add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=(30, 41, 55)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    _pptx_set_text(frame.paragraphs[0], text, size=size, bold=bold, color=color)
    return box


def _pptx_add_bullet_panel(slide, left, top, width, height, title, items, *, accent=(16, 36, 62)):
    panel = _pptx_add_panel(slide, left, top, width, height)
    frame = panel.text_frame
    frame.clear()
    frame.word_wrap = True
    _pptx_set_text(frame.paragraphs[0], title, size=18, bold=True, color=accent)
    for item in items:
        p = frame.add_paragraph()
        _pptx_set_text(p, f"• {item}", size=13, color=(71, 84, 103))
    return panel


def _pptx_static_path(static_url):
    value = str(static_url or "")
    if "/static/" in value:
        relative = value.split("/static/", 1)[1].replace("/", "\\")
        return Path("static") / Path(relative)
    return None


def _pptx_add_image(slide, static_url, left, top, width, height):
    local_path = _pptx_static_path(static_url)
    if not local_path or not local_path.exists():
        return False
    # python-pptx relies on Pillow for most image formats and can fail on SVG
    # or on files that are not valid raster images. In that case we gracefully
    # skip the image and let the caller render a textual fallback instead.
    if local_path.suffix.lower() == ".svg":
        return False
    try:
        slide.shapes.add_picture(str(local_path), left, top, width=width, height=height)
        return True
    except (UnidentifiedImageError, OSError, ValueError):
        return False


def _build_user_handbook_powerpoint(context):
    prs = Presentation()
    prs.slide_width = Mm(210)
    prs.slide_height = Mm(297)
    blank = prs.slide_layouts[6]

    navy = (16, 36, 62)
    blue = (31, 74, 114)
    slate = (71, 84, 103)
    light = (244, 247, 250)
    white = (255, 255, 255)

    def add_background(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)

    slide = prs.slides.add_slide(blank)
    add_background(slide, navy)
    _pptx_add_textbox(slide, Mm(14), Mm(16), Mm(120), Mm(16), context["closing_brand_name"], size=22, bold=True, color=white)
    _pptx_add_textbox(slide, Mm(14), Mm(42), Mm(126), Mm(22), "دليل المستخدم الفردي", size=28, bold=True, color=white)
    _pptx_add_textbox(slide, Mm(14), Mm(68), Mm(126), Mm(20), "تنزيل PowerPoint أصلي للدليل", size=16, color=(225, 232, 240))
    _pptx_add_bullet_panel(
        slide,
        Mm(14),
        Mm(98),
        Mm(118),
        Mm(72),
        "بيانات الإصدار",
        [
            f"المستخدم: {context['manual_target_name']}",
            f"اسم المستخدم: {context['manual_target_username']}",
            f"الصفة: {context['manual_target_position']}",
            f"تاريخ الإصدار: {context['manual_now_date']}",
        ],
        accent=white,
    )
    _pptx_add_bullet_panel(
        slide,
        Mm(138),
        Mm(98),
        Mm(58),
        Mm(72),
        "إحصاءات",
        [
            f"الشاشات: {context['manual_screen_total']}",
            f"الصفحات: {len(context['manual_content_pages'])}",
            f"العناصر: {context['manual_button_total']}",
        ],
        accent=white,
    )
    _pptx_add_textbox(slide, Mm(14), Mm(250), Mm(182), Mm(20), "ملف PowerPoint مولّد من نفس محتوى الدليل الحالي، ويمكن تعديله أو تصديره لاحقًا.", size=12, color=(230, 236, 241))

    slide = prs.slides.add_slide(blank)
    add_background(slide, light)
    _pptx_add_textbox(slide, Mm(12), Mm(12), Mm(186), Mm(16), "مقدمة الدليل", size=24, bold=True, color=navy)
    _pptx_add_bullet_panel(
        slide,
        Mm(12),
        Mm(34),
        Mm(90),
        Mm(88),
        "نطاق الاستخدام",
        [
            "نسخة فردية مخصصة للمستخدم الحالي.",
            "مناسبة للعرض والتصدير والطباعة.",
            "تشمل المراحل والشاشات والأخطاء والتنبيهات.",
        ],
        accent=navy,
    )
    _pptx_add_bullet_panel(
        slide,
        Mm(108),
        Mm(34),
        Mm(90),
        Mm(88),
        "مرجع سريع",
        [
            f"الأقسام: {' / '.join(context['manual_user_departments'])}",
            f"رابط النظام: {context['manual_target_system_url']}",
            f"رمز النظام: {context['manual_target_system_code']}",
        ],
        accent=blue,
    )

    for toc_page in context["manual_toc_pages"]:
        slide = prs.slides.add_slide(blank)
        add_background(slide, white)
        _pptx_add_textbox(slide, Mm(12), Mm(12), Mm(186), Mm(16), "فهرس الدليل", size=24, bold=True, color=navy)
        items = [f"{entry['page']:02d} — {entry['kind']} — {entry['title']}" for entry in toc_page["entries"]]
        _pptx_add_bullet_panel(slide, Mm(12), Mm(34), Mm(186), Mm(220), f"صفحة الفهرس {toc_page['toc_page_number']}", items, accent=blue)

    for page in context["manual_content_pages"]:
        slide = prs.slides.add_slide(blank)
        add_background(slide, white)
        _pptx_add_textbox(slide, Mm(12), Mm(10), Mm(186), Mm(12), page.get("eyebrow", ""), size=12, bold=True, color=blue)
        _pptx_add_textbox(slide, Mm(12), Mm(22), Mm(186), Mm(16), page["title"], size=22, bold=True, color=navy)
        _pptx_add_textbox(slide, Mm(12), Mm(38), Mm(186), Mm(24), page.get("intro", ""), size=12, color=slate)

        if page.get("layout") == "visual":
            image_added = _pptx_add_image(slide, page.get("image_path"), Mm(12), Mm(68), Mm(186), Mm(165))
            if not image_added:
                _pptx_add_bullet_panel(slide, Mm(12), Mm(68), Mm(186), Mm(110), page.get("image_title", "لقطة الشاشة"), [page.get("image_caption", "لا توجد صورة متاحة لهذه الشريحة.")], accent=navy)
            else:
                _pptx_add_textbox(slide, Mm(12), Mm(238), Mm(186), Mm(24), page.get("image_caption", ""), size=11, color=slate)
            continue

        if page["type"] == "workflow":
            main_items = [f"{item['index']}. {item['text']}" for item in page["items"]]
        else:
            main_items = [
                f"{item['index']}. {item['label']} | المكان: {item['location']} | الوظيفة: {item['purpose']}"
                for item in page["items"]
            ]

        side_items = []
        if page.get("used_by"):
            side_items.append(f"المستخدم: {page['used_by']}")
        if page.get("path"):
            side_items.append(f"المسار: {page['path']}")
        if page.get("continued"):
            side_items.append("هذه الشريحة تمثل استمرارًا منطقيًا لنفس القسم.")

        _pptx_add_bullet_panel(slide, Mm(12), Mm(68), Mm(122), Mm(185), "المحتوى", main_items, accent=navy)
        _pptx_add_bullet_panel(slide, Mm(138), Mm(68), Mm(60), Mm(88), "ملاحظات", side_items or ["لا توجد ملاحظات إضافية."], accent=blue)
        if page.get("show_image"):
            _pptx_add_image(slide, page.get("image_path"), Mm(138), Mm(162), Mm(60), Mm(68))

    slide = prs.slides.add_slide(blank)
    add_background(slide, light)
    _pptx_add_textbox(slide, Mm(12), Mm(12), Mm(186), Mm(16), "تنبيهات وتشغيل آمن", size=24, bold=True, color=navy)
    _pptx_add_bullet_panel(
        slide,
        Mm(12),
        Mm(34),
        Mm(186),
        Mm(150),
        "تعليمات مهمة",
        [
            "تحقق من الحساب والرابط قبل البدء.",
            "راجع الحقول الإلزامية قبل الحفظ أو التصدير.",
            "اختلاف الصلاحيات قد يغيّب بعض الشاشات أو الأزرار.",
            "استخدم النسخة المخصصة للإخراج النهائي.",
        ],
        accent=blue,
    )

    slide = prs.slides.add_slide(blank)
    add_background(slide, navy)
    _pptx_add_textbox(slide, Mm(12), Mm(18), Mm(186), Mm(16), "خاتمة الدليل", size=24, bold=True, color=white)
    _pptx_add_bullet_panel(
        slide,
        Mm(12),
        Mm(46),
        Mm(90),
        Mm(90),
        "جهة الدعم",
        [
            context["closing_developer_name"],
            context["closing_role"],
            context["closing_developer_phone"],
            context["closing_developer_email"],
        ],
        accent=white,
    )
    _pptx_add_bullet_panel(
        slide,
        Mm(108),
        Mm(46),
        Mm(90),
        Mm(90),
        "مزايا الملف",
        [
            "PowerPoint أصلي قابل للتعديل.",
            "مستقل عن نسخة الطباعة HTML.",
            "مناسب للتقديم أو التصدير إلى PDF.",
        ],
        accent=white,
    )

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def _build_user_handbook_powerpoint_from_rendered_pages(request, context):
    edge_path = _find_edge_executable()
    if not edge_path:
        return _build_user_handbook_powerpoint(context)

    prs = Presentation()
    prs.slide_width = Mm(210)
    prs.slide_height = Mm(297)
    blank = prs.slide_layouts[6]

    total_pages = _manual_total_render_pages(context)
    base_query = request.GET.copy()
    base_query.pop("export_page", None)
    base_query["export_mode"] = "powerpoint"
    signed_payload = signing.dumps(
        {key: base_query.getlist(key) for key in base_query.keys()},
        salt=MANUAL_EXPORT_SIGNING_SALT,
        compress=True,
    )

    try:
        with tempfile.TemporaryDirectory(prefix="manual-pptx-") as temp_dir:
            temp_path = Path(temp_dir)
            export_print_url = _reverse_manual_url("manuals:user_handbook_print_export", "user_handbook_print_export")
            for page_number in range(1, total_pages + 1):
                screenshot_path = temp_path / f"page-{page_number:03d}.png"
                url = request.build_absolute_uri(
                    f"{export_print_url}?token={signed_payload}&export_page={page_number}"
                )
                subprocess.run(
                    [
                        str(edge_path),
                        "--headless=new",
                        "--disable-gpu",
                        "--hide-scrollbars",
                        "--default-background-color=ffffff",
                        "--force-device-scale-factor=1",
                        "--window-size=794,1123",
                        "--virtual-time-budget=5000",
                        f"--screenshot={screenshot_path}",
                        url,
                    ],
                    check=True,
                    timeout=45,
                )

                if not screenshot_path.exists():
                    continue

                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(str(screenshot_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    except (subprocess.SubprocessError, OSError, ValueError):
        return _build_user_handbook_powerpoint(context)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


class ManualsHomeView(LoginRequiredMixin, TemplateView):
    template_name = "manuals/home.html"

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context.update(
            build_manuals_context(
                query=(self.request.GET.get("q") or "").strip().lower(),
                selected_group=(self.request.GET.get("group") or "").strip(),
            )
        )
        context.update(_build_closing_page_context(self.request))
        return context


class ManualsHandbookView(LoginRequiredMixin, TemplateView):
    template_name = "manuals/handbook.html"

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context.update(
            build_manuals_context(
                query=(self.request.GET.get("q") or "").strip().lower(),
                selected_group=(self.request.GET.get("group") or "").strip(),
            )
        )
        context.update(_build_closing_page_context(self.request))
        return context


class ManualsUserGuideSelectView(LoginRequiredMixin, TemplateView):
    template_name = "manuals/user_select.html"

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        selected_user, users = _selected_manual_user(self.request)
        context["manual_target_users"] = users
        context["manual_selected_user"] = selected_user
        context["manual_selected_password"] = (self.request.GET.get("password") or "").strip()
        return context


class ManualsUserHandbookView(LoginRequiredMixin, TemplateView):
    template_name = "manuals/user_handbook.html"

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context.update(_build_user_handbook_page_context(self.request, print_mode=False))
        return context


class ManualsUserHandbookPrintView(ManualsUserHandbookView):
    template_name = "manuals/user_handbook_print.html"

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context["manual_print_mode"] = True
        context["manual_print_url"] = self.request.build_absolute_uri()
        return context


class ManualsUserHandbookPrintExportView(TemplateView):
    template_name = "manuals/user_handbook_print.html"

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        token = (self.request.GET.get("token") or "").strip()
        payload = signing.loads(token, salt=MANUAL_EXPORT_SIGNING_SALT, max_age=900)
        params = {}
        for key, value in payload.items():
            if isinstance(value, list):
                params[key] = value[0] if value else ""
            else:
                params[key] = value
        params["export_page"] = (self.request.GET.get("export_page") or "").strip()
        params["export_mode"] = "powerpoint"
        context.update(_build_user_handbook_page_context(self.request, params=params, print_mode=True))
        context["manual_print_url"] = self.request.build_absolute_uri()
        return context


class ManualsMarkdownDownloadView(LoginRequiredMixin, TemplateView):
    def get(self, request, *args, **kwargs):
        file_path = Path("docs/user-guide.md")
        return FileResponse(
            file_path.open("rb"),
            as_attachment=True,
            filename="user-guide.md",
            content_type="text/markdown; charset=utf-8",
        )


class ManualsUserHandbookPowerPointDownloadView(ManualsUserHandbookView):
    def get(self, request, *args, **kwargs):
        context = self.get_context_data(**kwargs)
        presentation = _build_user_handbook_powerpoint_from_rendered_pages(request, context)
        filename = f"user-handbook-{context['manual_target_username']}.pptx"
        response = HttpResponse(
            presentation.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response
